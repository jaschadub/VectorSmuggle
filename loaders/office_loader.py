"""Office document loader for DOCX, XLSX, PPTX files."""

import logging
from pathlib import Path

from langchain.schema import Document


class OfficeLoader:
    """Loader for Microsoft Office documents (DOCX, XLSX, PPTX)."""

    def __init__(self, file_path: str, logger: logging.Logger | None = None):
        """Initialize the office document loader.

        Args:
            file_path: Path to the office document
            logger: Optional logger instance
        """
        self.file_path = Path(file_path)
        self.logger = logger or logging.getLogger(__name__)

        if not self.file_path.exists():
            raise FileNotFoundError(f"File not found: {self.file_path}")

        self.file_type = self.file_path.suffix.lower()

        # Validate supported formats
        supported_formats = ['.docx', '.xlsx', '.pptx', '.doc', '.xls', '.ppt']
        if self.file_type not in supported_formats:
            raise ValueError(f"Unsupported office format: {self.file_type}")

    def load(self) -> list[Document]:
        """Load the office document and return Document objects.

        Returns:
            List of Document objects with extracted content

        Raises:
            ImportError: If required libraries are not installed
            Exception: If document loading fails
        """
        try:
            if self.file_type in ['.docx', '.doc']:
                return self._load_word_document()
            elif self.file_type in ['.xlsx', '.xls']:
                return self._load_excel_document()
            elif self.file_type in ['.pptx', '.ppt']:
                return self._load_powerpoint_document()
            else:
                raise ValueError(f"Unsupported file type: {self.file_type}")

        except ImportError as e:
            self.logger.error(f"Missing required library for {self.file_type}: {e}")
            raise ImportError(f"Missing required library for {self.file_type}") from e
        except Exception as e:
            self.logger.error(f"Failed to load {self.file_path}: {e}")
            raise

    def _load_word_document(self) -> list[Document]:
        """Load Word document (.docx/.doc).

        Returns:
            List of Document objects
        """
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("python-docx is required for DOCX files. Install with: pip install python-docx")

        try:
            doc = DocxDocument(self.file_path)

            # Extract text content
            content_parts = []

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_parts.append(paragraph.text.strip())

            # Extract tables
            for table in doc.tables:
                table_text = self._extract_table_text(table)
                if table_text:
                    content_parts.append(f"[TABLE]\n{table_text}")

            # Combine content
            full_content = "\n\n".join(content_parts)

            # Extract metadata
            metadata = self._extract_word_metadata(doc)
            metadata.update({
                'source': str(self.file_path.resolve()),
                'file_type': self.file_type,
                'loader': 'OfficeLoader'
            })

            self.logger.info(f"Extracted {len(content_parts)} content sections from Word document")

            return [Document(page_content=full_content, metadata=metadata)]

        except Exception as e:
            self.logger.error(f"Failed to load Word document: {e}")
            raise

    def _load_excel_document(self) -> list[Document]:
        """Load Excel document (.xlsx/.xls).

        Returns:
            List of Document objects (one per worksheet)
        """
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for Excel files. Install with: pip install openpyxl")

        try:
            workbook = openpyxl.load_workbook(self.file_path, data_only=True)
            documents = []

            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]

                # Extract data from worksheet
                content_parts = []

                # Get all rows with data
                for row in worksheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):
                        content_parts.append("\t".join(row_data))

                if content_parts:
                    sheet_content = "\n".join(content_parts)

                    metadata = {
                        'source': str(self.file_path.resolve()),
                        'file_type': self.file_type,
                        'loader': 'OfficeLoader',
                        'sheet_name': sheet_name,
                        'rows_count': len(content_parts),
                        'columns_count': worksheet.max_column
                    }

                    documents.append(Document(
                        page_content=f"[WORKSHEET: {sheet_name}]\n{sheet_content}",
                        metadata=metadata
                    ))

            self.logger.info(f"Extracted {len(documents)} worksheets from Excel document")
            return documents

        except Exception as e:
            self.logger.error(f"Failed to load Excel document: {e}")
            raise

    def _load_powerpoint_document(self) -> list[Document]:
        """Load PowerPoint document (.pptx/.ppt).

        Returns:
            List of Document objects (one per slide)
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PowerPoint files. Install with: pip install python-pptx")

        try:
            presentation = Presentation(self.file_path)
            documents = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                content_parts = []

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text.strip())

                    # Extract table content if present
                    if shape.has_table:
                        table_text = self._extract_pptx_table_text(shape.table)
                        if table_text:
                            content_parts.append(f"[TABLE]\n{table_text}")

                if content_parts:
                    slide_content = "\n\n".join(content_parts)

                    metadata = {
                        'source': str(self.file_path.resolve()),
                        'file_type': self.file_type,
                        'loader': 'OfficeLoader',
                        'slide_number': slide_num,
                        'total_slides': len(presentation.slides)
                    }

                    documents.append(Document(
                        page_content=f"[SLIDE {slide_num}]\n{slide_content}",
                        metadata=metadata
                    ))

            self.logger.info(f"Extracted {len(documents)} slides from PowerPoint document")
            return documents

        except Exception as e:
            self.logger.error(f"Failed to load PowerPoint document: {e}")
            raise

    def _extract_table_text(self, table) -> str:
        """Extract text from Word table.

        Args:
            table: Word table object

        Returns:
            Formatted table text
        """
        table_text = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                row_text.append(cell.text.strip())
            table_text.append("\t".join(row_text))
        return "\n".join(table_text)

    def _extract_pptx_table_text(self, table) -> str:
        """Extract text from PowerPoint table.

        Args:
            table: PowerPoint table object

        Returns:
            Formatted table text
        """
        table_text = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                row_text.append(cell.text.strip())
            table_text.append("\t".join(row_text))
        return "\n".join(table_text)

    def _extract_word_metadata(self, doc) -> dict:
        """Extract metadata from Word document.

        Args:
            doc: Word document object

        Returns:
            Dictionary of metadata
        """
        metadata = {}

        try:
            core_props = doc.core_properties

            if core_props.author:
                metadata['author'] = core_props.author
            if core_props.title:
                metadata['title'] = core_props.title
            if core_props.subject:
                metadata['subject'] = core_props.subject
            if core_props.created:
                metadata['created'] = core_props.created.isoformat()
            if core_props.modified:
                metadata['modified'] = core_props.modified.isoformat()
            if core_props.last_modified_by:
                metadata['last_modified_by'] = core_props.last_modified_by
            if core_props.revision:
                metadata['revision'] = core_props.revision

        except Exception as e:
            self.logger.warning(f"Failed to extract Word metadata: {e}")

        return metadata
